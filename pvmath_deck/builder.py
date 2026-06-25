"""python-pptx slide builder for PVMath branded decks."""

from __future__ import annotations

from datetime import date
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from pvmath_deck.content import DISCLAIMER
from pvmath_deck.theme import BRAND, Brand


def _rgb(t: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*t)


class DeckBuilder:
    """Fluent builder for 16:9 PVMath decks."""

    def __init__(
        self,
        brand: Brand = BRAND,
        presenter: str = "[Presenter name]",
        deck_date: str | None = None,
    ) -> None:
        self.b = brand
        self.presenter = presenter
        self.deck_date = deck_date or date.today().strftime("%d %B %Y")
        self.prs = Presentation()
        self.prs.slide_width = brand.slide_w
        self.prs.slide_height = brand.slide_h
        self._slide_num = 0

    def save(self, path: str) -> None:
        self.prs.save(path)

    # ── Primitives ────────────────────────────────────────────────────────

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bar(self, slide, height=Inches(0.08)):
        sh = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, self.b.slide_w, height
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(self.b.green)
        sh.line.fill.background()

    def _footer(self, slide, left_text: str = ""):
        y = self.b.slide_h - Inches(0.38)
        if left_text:
            box = slide.shapes.add_textbox(
                self.b.margin, y, Inches(6), Inches(0.3)
            )
            tf = box.text_frame
            tf.text = left_text
            p = tf.paragraphs[0]
            p.font.size = Pt(9)
            p.font.name = self.b.font_body
            p.font.color.rgb = _rgb(self.b.muted)
        box = slide.shapes.add_textbox(
            self.b.slide_w - Inches(3.2), y, Inches(3), Inches(0.3)
        )
        tf = box.text_frame
        tf.text = self.b.website.replace("https://", "")
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(9)
        p.font.name = self.b.font_body
        p.font.color.rgb = _rgb(self.b.muted)

    def _title_block(self, slide, title: str, subtitle: str = ""):
        box = slide.shapes.add_textbox(
            self.b.margin, Inches(0.45), self.b.slide_w - Inches(1.1), Inches(1.2)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.b.font_title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = _rgb(self.b.text)
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.name = self.b.font_body
            p2.font.size = Pt(16)
            p2.font.color.rgb = _rgb(self.b.muted)
            p2.space_before = Pt(6)

    def _bullets(
        self,
        slide,
        items: Sequence[str],
        left=Inches,
        top=Inches,
        width=Inches,
        height=Inches,
        size: int = 18,
    ):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = self.b.font_body
            p.font.size = Pt(size)
            p.font.color.rgb = _rgb(self.b.text)
            p.space_after = Pt(8)

    def _placeholder(self, slide, label: str, left, top, width, height):
        sh = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(self.b.placeholder_bg)
        sh.line.color.rgb = _rgb(self.b.border)
        tf = sh.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"[INSERT IMAGE]\n{label}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.name = self.b.font_body
        p.font.color.rgb = _rgb(self.b.muted)

    def _notes(
        self,
        slide,
        body: str,
        duration_sec: int = 60,
        transition: str = "",
        qa: Sequence[tuple[str, str]] = (),
    ):
        self._slide_num += 1
        parts = [f"Slide {self._slide_num}", f"Duration: ~{duration_sec} sec"]
        if transition:
            parts.append(f"Transition: {transition}")
        parts.append("")
        parts.append(body)
        if qa:
            parts.append("")
            parts.append("— Anticipated Q&A —")
            for q, a in qa:
                parts.append(f"Q: {q}")
                parts.append(f"A: {a}")
        notes = slide.notes_slide.notes_text_frame
        notes.text = "\n".join(parts)

    def _logo_mark(self, slide, left, top, size=Inches(0.65)):
        # Rounded square + PV text (matches assets/logo.svg)
        sh = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, size, size
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(self.b.green_dk)
        sh.line.fill.background()
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, size, size * 0.28
        )
        band.fill.solid()
        band.fill.fore_color.rgb = _rgb(self.b.green)
        band.line.fill.background()
        tb = slide.shapes.add_textbox(
            left, top + size * 0.35, size, size * 0.5
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "PV"
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = _rgb(self.b.white)

    # ── Slide types ───────────────────────────────────────────────────────

    def cover(self, title: str, subtitle: str, notes: str, duration: int = 30):
        slide = self._blank()
        self._bar(slide)
        self._logo_mark(slide, self.b.margin, Inches(1.0), Inches(0.9))
        box = slide.shapes.add_textbox(
            self.b.margin + Inches(1.05), Inches(0.95), Inches(4), Inches(0.6)
        )
        p = box.text_frame.paragraphs[0]
        p.text = self.b.name
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = _rgb(self.b.text)

        box = slide.shapes.add_textbox(
            self.b.margin, Inches(2.0), Inches(10), Inches(2.2)
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = _rgb(self.b.text)
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = _rgb(self.b.muted)
            p2.space_before = Pt(10)

        meta = slide.shapes.add_textbox(
            self.b.margin, Inches(5.8), Inches(8), Inches(0.8)
        )
        tf = meta.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.presenter}  ·  {self.deck_date}"
        p.font.size = Pt(14)
        p.font.color.rgb = _rgb(self.b.muted)

        disc = slide.shapes.add_textbox(
            self.b.margin, Inches(6.35), Inches(11), Inches(0.5)
        )
        tf = disc.text_frame
        p = tf.paragraphs[0]
        p.text = DISCLAIMER
        p.font.size = Pt(9)
        p.font.italic = True
        p.font.color.rgb = _rgb(self.b.muted)

        self._footer(slide)
        self._notes(slide, notes, duration, transition="Open with customer context.")

    def agenda(self, items: Sequence[str], notes: str, duration: int = 45):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, "Agenda")
        self._bullets(
            slide,
            [f"{i + 1}. {t}" for i, t in enumerate(items)],
            self.b.margin,
            Inches(1.5),
            Inches(11),
            Inches(5),
            size=20,
        )
        self._footer(slide)
        self._notes(slide, notes, duration)

    def bullets_slide(
        self,
        title: str,
        items: Sequence[str],
        notes: str,
        duration: int = 60,
        subtitle: str = "",
        visual: str = "",
        qa: Sequence[tuple[str, str]] = (),
    ):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, title, subtitle)
        w = Inches(6.2) if visual else Inches(11.5)
        self._bullets(slide, items, self.b.margin, Inches(1.55), w, Inches(5.2))
        if visual:
            self._placeholder(
                slide,
                visual,
                Inches(7.0),
                Inches(1.55),
                Inches(5.8),
                Inches(4.8),
            )
        self._footer(slide, DISCLAIMER[:60] + "…")
        self._notes(slide, notes, duration, qa=qa)

    def two_column(
        self,
        title: str,
        left_title: str,
        left_items: Sequence[str],
        right_title: str,
        right_items: Sequence[str],
        notes: str,
        duration: int = 75,
    ):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, title)
        lx, rx = self.b.margin, Inches(6.85)
        for x, ht, items in ((lx, left_title, left_items), (rx, right_title, right_items)):
            box = slide.shapes.add_textbox(x, Inches(1.45), Inches(5.9), Inches(0.4))
            p = box.text_frame.paragraphs[0]
            p.text = ht
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = _rgb(self.b.green_dk)
            self._bullets(slide, items, x, Inches(1.95), Inches(5.9), Inches(4.5), size=16)
        self._footer(slide)
        self._notes(slide, notes, duration)

    def table_slide(
        self,
        title: str,
        headers: Sequence[str],
        rows: Sequence[Sequence[str]],
        notes: str,
        duration: int = 75,
        subtitle: str = "",
    ):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, title, subtitle)
        nrows, ncols = len(rows) + 1, len(headers)
        tbl = slide.shapes.add_table(
            nrows, ncols, self.b.margin, Inches(1.55), Inches(12.2), Inches(0.5 * nrows)
        ).table
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(self.b.green_dk)
            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(12)
                p.font.color.rgb = _rgb(self.b.white)
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = _rgb(self.b.text)
        self._footer(slide)
        self._notes(slide, notes, duration)

    def module_card_slide(
        self,
        title: str,
        modules: Sequence[dict],
        notes: str,
        duration: int = 90,
    ):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, title)
        n = len(modules)
        gap = Inches(0.25)
        w = (Inches(12.2) - gap * (n - 1)) / n
        for i, m in enumerate(modules):
            left = self.b.margin + (w + gap) * i
            sh = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                left,
                Inches(1.55),
                w,
                Inches(4.9),
            )
            sh.fill.solid()
            sh.fill.fore_color.rgb = _rgb(self.b.green_lt)
            sh.line.color.rgb = _rgb(self.b.border)
            tf = sh.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.15)
            tf.margin_top = Inches(0.12)
            lines = [
                m["name"],
                m.get("tagline", ""),
                "",
                f"Outputs: {m.get('outputs', '')[:80]}…" if len(m.get("outputs", "")) > 80 else f"Outputs: {m.get('outputs', '')}",
                f"Typical: {m.get('time', '')}",
            ]
            for j, line in enumerate(lines):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16 if j == 0 else 11)
                p.font.bold = j == 0
                p.font.color.rgb = _rgb(self.b.green_dk if j == 0 else self.b.text)
        self._footer(slide)
        self._notes(slide, notes, duration)

    def demo_steps(
        self,
        title: str,
        steps: Sequence[dict],
        notes: str,
        duration: int = 120,
    ):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, title, "Live demo flow")
        y = Inches(1.5)
        for i, step in enumerate(steps[:4]):  # max 4 on one slide; deck may split
            box = slide.shapes.add_textbox(self.b.margin, y, Inches(5.5), Inches(1.1))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Step {i + 1} — {step['title']}"
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = _rgb(self.b.green_dk)
            for b in step.get("bullets", [])[:2]:
                p2 = tf.add_paragraph()
                p2.text = f"• {b}"
                p2.font.size = Pt(11)
                p2.font.color.rgb = _rgb(self.b.text)
            y += Inches(1.15)
        if steps:
            self._placeholder(
                slide,
                steps[0].get("visual", "Demo screenshot"),
                Inches(6.8),
                Inches(1.5),
                Inches(5.9),
                Inches(4.9),
            )
        self._footer(slide)
        self._notes(
            slide,
            notes + "\n\n" + "\n".join(
                f"Step {i+1}: {s.get('talk', '')}" for i, s in enumerate(steps)
            ),
            duration,
        )

    def architecture(self, rows: Sequence[tuple[str, str]], notes: str, duration: int = 75):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, "Technical architecture", "High level — Early Access")
        self._placeholder(
            slide,
            "Optional: architecture diagram export",
            Inches(7.2),
            Inches(1.5),
            Inches(5.5),
            Inches(4.8),
        )
        self._bullets(
            slide,
            [f"{a}: {b}" for a, b in rows],
            self.b.margin,
            Inches(1.55),
            Inches(6.3),
            Inches(4.8),
            size=15,
        )
        self._footer(slide)
        self._notes(slide, notes, duration)

    def roadmap(self, items: Sequence[tuple[str, str]], notes: str, duration: int = 60):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, "Product roadmap")
        self._bullets(
            slide,
            [f"{phase}: {detail}" for phase, detail in items],
            self.b.margin,
            Inches(1.55),
            Inches(11),
            Inches(4.5),
            size=18,
        )
        self._footer(slide)
        self._notes(slide, notes, duration)

    def pricing(self, notes: str, duration: int = 45):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, "Pricing", "Early Access — contact for activation")
        from pvmath_deck.content import PRICING

        self._bullets(
            slide,
            PRICING.split(" · "),
            self.b.margin,
            Inches(1.55),
            Inches(11),
            Inches(3),
            size=17,
        )
        box = slide.shapes.add_textbox(self.b.margin, Inches(4.2), Inches(11), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = f"Try free: {self.b.app_url}  ·  Subscribe / Enterprise: {self.b.email}"
        p.font.size = Pt(14)
        p.font.color.rgb = _rgb(self.b.green_dk)
        self._footer(slide)
        self._notes(slide, notes, duration)

    def qa_contact(self, notes: str, duration: int = 30):
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, "Questions?", "Thank you")
        lines = [
            self.b.website,
            self.b.email,
            self.b.linkedin,
            self.b.youtube,
            f"App: {self.b.app_url}",
        ]
        self._bullets(slide, lines, self.b.margin, Inches(2.2), Inches(10), Inches(3.5), size=20)
        self._logo_mark(slide, Inches(10.5), Inches(2.0), Inches(1.1))
        self._footer(slide)
        self._notes(slide, notes, duration, transition="Offer follow-up demo or pilot proposal.")

    def speaker_script_appendix(self, sections: dict[str, str]) -> None:
        """Hidden appendix slide with full demo script (presenter reference)."""
        slide = self._blank()
        self._bar(slide)
        self._title_block(slide, "Appendix — Speaker script", "Hide or skip in live delivery")
        text = []
        for heading, body in sections.items():
            text.append(f"▸ {heading}")
            text.append(body)
            text.append("")
        self._bullets(slide, text[:12], self.b.margin, Inches(1.45), Inches(12), Inches(5.5), size=10)
        self._footer(slide, "Presenter reference")
        self._notes(
            slide,
            "Full script for presenter — not shown to audience unless needed.",
            0,
        )
